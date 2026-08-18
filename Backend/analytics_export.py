"""
Turning a dashboard view into a file someone can send to somebody else.

Five formats, one shared summary: PDF, Word, PowerPoint, a Power BI-ready
workbook, and a Power BI connection file.

Charts arrive as PNGs captured from the browser rather than being redrawn here.
That is deliberate: the alternative is a second chart implementation in Python
which would slowly disagree with the one on screen, and an export that doesn't
match what the tester was looking at is worse than no export.

A .pbix is NOT produced. The format is proprietary and Microsoft publishes no
authoring API, so anything claiming to emit one would be guesswork. The honest
equivalent is a clean dataset plus a connection file, which is what an analyst
needs to build visuals against live, refreshable data.
"""
import base64
import io
import json
from datetime import datetime, timezone

RANGE_LABELS = {
    "1m": "last month", "3m": "last 3 months", "6m": "last 6 months",
    "1y": "last year", "all": "all time",
}
BASIS_LABELS = {"latest": "latest run per test layer", "all": "all runs"}


def _fmt(n):
    return f"{n:,}" if isinstance(n, (int, float)) and n is not None else "--"


def describe_scope(payload):
    """One line saying exactly what the reader is looking at."""
    summary = payload.get("summary", {})
    layers = payload.get("scope_layers") or []
    who = ", ".join(layers) if layers else "All test layers"
    return (f"{who}  |  {BASIS_LABELS.get(payload.get('basis'), payload.get('basis'))}"
            f"  |  {RANGE_LABELS.get(payload.get('range'), payload.get('range'))}"
            f"  |  {summary.get('runs_covered', 0)} runs, "
            f"{summary.get('layers_covered', 0)} test layers")


def headline_rows(payload):
    """The KPI block as label/value/detail triples - shared by every format."""
    s = payload.get("summary", {})
    checks, rows, nulls = s.get("checks", {}), s.get("rows", {}), s.get("nulls", {})
    total = checks.get("total") or 0
    pass_pct = round(checks.get("pass", 0) / total * 100) if total else 0
    quality = rows.get("quality_pct")
    return [
        ("Checks passed", f"{checks.get('pass', 0)} of {total}", f"{pass_pct}% of checks run"),
        ("Checks failed", _fmt(checks.get("fail", 0)),
         f"{checks.get('error', 0)} errored (could not run)"),
        ("Row-level quality", "--" if quality is None else f"{quality}%",
         f"{_fmt(rows.get('clean'))} clean of {_fmt(rows.get('examined'))} rows examined"),
        ("Null values found", _fmt(nulls.get("violations")),
         f"across {nulls.get('checks', 0)} null checks"),
    ]


def caveats(payload):
    """
    The facts that stop the numbers being misread. Carried into every export,
    because a document travels away from the screen that explained them.
    """
    s = payload.get("summary", {})
    notes = []
    excluded = s.get("rows", {}).get("excluded_results") or 0
    if excluded:
        notes.append(
            f"{excluded} checks report no row counts (comparisons, freshness, categorical) and are "
            "excluded from row-level quality rather than counted as zero.")
    if s.get("orphaned_runs"):
        notes.append(
            f"Includes {s['orphaned_runs']} runs whose test layer has since been deleted, shown as "
            "'(deleted test layer)'." if s.get("orphaned_runs_included") else
            f"{s['orphaned_runs']} runs from deleted test layers are excluded while a layer filter "
            "is applied.")
    notes.append(
        "Violation counts are exactly what each check reported. A cross-table check can legitimately "
        "report more violations than the table has rows, so row-level quality clamps per check.")
    return notes


def _decode_charts(charts):
    """[{title, image: 'data:image/png;base64,...'}] -> [(title, bytes)]."""
    out = []
    for c in charts or []:
        raw = c.get("image") or ""
        if "," in raw:
            raw = raw.split(",", 1)[1]
        if not raw:
            continue
        try:
            out.append((c.get("title") or "Chart", base64.b64decode(raw)))
        except Exception:
            continue  # one chart that failed to capture must not sink the export
    return out


# --- Word -------------------------------------------------------------------

def to_docx(payload, charts):
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor

    doc = Document()
    doc.add_heading("Data Quality Report", level=0)
    doc.add_paragraph(describe_scope(payload))
    doc.add_paragraph(f"Generated {datetime.now(timezone.utc).strftime('%d %b %Y %H:%M UTC')}")

    def table(headers, rows):
        t = doc.add_table(rows=1, cols=len(headers))
        t.style = "Light Grid Accent 1"
        for cell, text in zip(t.rows[0].cells, headers):
            cell.text = text
        for row in rows:
            cells = t.add_row().cells
            for cell, value in zip(cells, row):
                cell.text = str(value)
        return t

    doc.add_heading("Headline", level=1)
    table(["Measure", "Value", "Detail"], [list(r) for r in headline_rows(payload)])

    decoded = _decode_charts(charts)
    if decoded:
        doc.add_heading("Charts", level=1)
        for title, image in decoded:
            doc.add_paragraph(title, style="Heading 3")
            doc.add_picture(io.BytesIO(image), width=Inches(6.0))

    by_type = payload.get("by_validation_type") or []
    if by_type:
        doc.add_heading("By quality dimension", level=1)
        table(["Dimension", "Passed", "Failed", "Errored", "Violating rows"],
              [[r.get("type"), _fmt(r.get("pass")), _fmt(r.get("fail")), _fmt(r.get("error")),
                "--" if r.get("violation_pct") is None else f"{r['violation_pct']}%"]
               for r in by_type])

    offenders = payload.get("worst_offenders") or []
    if offenders:
        doc.add_heading("Where the violations are", level=1)
        table(["Check", "Test layer", "Dimension", "Violations", "Status"],
              [[o.get("test_name"), o.get("mapping_name"), o.get("validation_type"),
                _fmt(o.get("violations")), o.get("status")] for o in offenders])

    doc.add_heading("How to read this", level=1)
    for note in caveats(payload):
        p = doc.add_paragraph(note, style="List Bullet")
        p.runs[0].font.size = Pt(9)
        p.runs[0].font.color.rgb = RGBColor(0x52, 0x52, 0x52)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# --- PowerPoint -------------------------------------------------------------

def to_pptx(payload, charts):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)   # 16:9
    brand = RGBColor(0x00, 0x16, 0x89)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Data Quality Report"
    title_slide.placeholders[1].text = describe_scope(payload)

    # One number per line, no chartjunk - this slide is read from the back row.
    kpi = prs.slides.add_slide(prs.slide_layouts[5])
    kpi.shapes.title.text = "Headline"
    top = Inches(1.8)
    for label, value, detail in headline_rows(payload):
        tf = kpi.shapes.add_textbox(Inches(0.8), top, Inches(11.5), Inches(1.0)).text_frame
        tf.text = f"{label}: {value}"
        run = tf.paragraphs[0].runs[0]
        run.font.size, run.font.bold, run.font.color.rgb = Pt(24), True, brand
        p = tf.add_paragraph()
        p.text = detail
        p.runs[0].font.size = Pt(13)
        top += Inches(1.15)

    for title, image in _decode_charts(charts):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        slide.shapes.add_picture(io.BytesIO(image), Inches(1.4), Inches(1.6), width=Inches(10.5))

    offenders = payload.get("worst_offenders") or []
    if offenders:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Where the violations are"
        rows = min(len(offenders), 10) + 1
        table = slide.shapes.add_table(
            rows, 4, Inches(0.7), Inches(1.6), Inches(11.9), Inches(0.4 * rows)).table
        for i, head in enumerate(("Check", "Test layer", "Dimension", "Violations")):
            table.cell(0, i).text = head
        for r, o in enumerate(offenders[:10], start=1):
            table.cell(r, 0).text = str(o.get("test_name"))[:40]
            table.cell(r, 1).text = str(o.get("mapping_name"))[:28]
            table.cell(r, 2).text = str(o.get("validation_type"))
            table.cell(r, 3).text = _fmt(o.get("violations"))

    notes = prs.slides.add_slide(prs.slide_layouts[5])
    notes.shapes.title.text = "How to read this"
    tf = notes.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.5)).text_frame
    tf.word_wrap = True
    for i, note in enumerate(caveats(payload)):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"- {note}"
        p.runs[0].font.size = Pt(14)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- PDF --------------------------------------------------------------------

def to_pdf(payload, charts):
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib.utils import ImageReader
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak)

    styles = getSampleStyleSheet()
    small = ParagraphStyle("small", parent=styles["Normal"], fontSize=8,
                           textColor=colors.HexColor("#525252"), leading=11)
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, title="Data Quality Report",
                            leftMargin=18 * mm, rightMargin=18 * mm,
                            topMargin=16 * mm, bottomMargin=16 * mm)
    width = doc.width

    def styled(data):
        t = Table(data, hAlign="LEFT", colWidths=[width / len(data[0])] * len(data[0]))
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#f1f5f9")),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#e2e8f0")),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#fafafa")]),
        ]))
        return t

    flow = [
        Paragraph("Data Quality Report", styles["Title"]),
        Paragraph(describe_scope(payload), small),
        Paragraph(f"Generated {datetime.now(timezone.utc).strftime('%d %b %Y %H:%M UTC')}", small),
        Spacer(1, 8 * mm),
        Paragraph("Headline", styles["Heading2"]),
        styled([["Measure", "Value", "Detail"]] + [list(r) for r in headline_rows(payload)]),
        Spacer(1, 6 * mm),
    ]

    for title, image in _decode_charts(charts):
        iw, ih = ImageReader(io.BytesIO(image)).getSize()
        flow += [Paragraph(title, styles["Heading3"]),
                 Image(io.BytesIO(image), width=width, height=width * ih / iw),
                 Spacer(1, 5 * mm)]

    by_type = payload.get("by_validation_type") or []
    if by_type:
        flow += [PageBreak(), Paragraph("By quality dimension", styles["Heading2"]),
                 styled([["Dimension", "Passed", "Failed", "Errored", "Violating rows"]]
                        + [[str(r.get("type")), _fmt(r.get("pass")), _fmt(r.get("fail")),
                            _fmt(r.get("error")),
                            "--" if r.get("violation_pct") is None else f"{r['violation_pct']}%"]
                           for r in by_type]),
                 Spacer(1, 6 * mm)]

    offenders = payload.get("worst_offenders") or []
    if offenders:
        flow += [Paragraph("Where the violations are", styles["Heading2"]),
                 styled([["Check", "Test layer", "Dimension", "Violations", "Status"]]
                        + [[str(o.get("test_name"))[:38], str(o.get("mapping_name"))[:24],
                            str(o.get("validation_type")), _fmt(o.get("violations")),
                            str(o.get("status"))] for o in offenders]),
                 Spacer(1, 6 * mm)]

    flow.append(Paragraph("How to read this", styles["Heading2"]))
    for note in caveats(payload):
        flow += [Paragraph(f"&bull; {note}", small), Spacer(1, 2 * mm)]

    doc.build(flow)
    return buf.getvalue()


# --- Power BI ---------------------------------------------------------------

def to_xlsx(payload, _charts=None):
    """
    A flat, Power BI-ready workbook: one row per check on the first sheet, so it
    imports as a table needing no reshaping. The aggregates get their own sheets
    for anyone who wants the totals without recomputing them.
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    wb.remove(wb.active)
    head_font = Font(bold=True, color="FFFFFF")
    head_fill = PatternFill("solid", fgColor="001689")

    def sheet(title, headers, rows):
        ws = wb.create_sheet(title)
        ws.append(headers)
        for cell in ws[1]:
            cell.font, cell.fill = head_font, head_fill
        for row in rows:
            ws.append(row)
        for i, header in enumerate(headers, start=1):
            widest = max([len(str(header))] + [len(str(r[i - 1])) for r in rows] + [0])
            ws.column_dimensions[get_column_letter(i)].width = min(max(widest + 2, 10), 46)
        ws.freeze_panes = "A2"

    sheet("Checks",
          ["Test layer", "Check", "Target", "Dimension", "Status", "Violations",
           "Rows examined", "Run started"],
          [[r.get("mapping_name"), r.get("test_name"), r.get("rule_target"),
            r.get("validation_type"), r.get("status"), r.get("violations"),
            r.get("total_rows"), r.get("run_started_at")]
           for r in (payload.get("checks_detail") or [])])

    sheet("By dimension",
          ["Dimension", "Passed", "Failed", "Errored", "Violations", "Rows", "Violation %"],
          [[r.get("type"), r.get("pass"), r.get("fail"), r.get("error"),
            r.get("violations"), r.get("rows"), r.get("violation_pct")]
           for r in (payload.get("by_validation_type") or [])])

    sheet("By test layer",
          ["Test layer", "Passed", "Failed", "Errored", "Row quality %", "Last run"],
          [[r.get("name"), r.get("pass"), r.get("fail"), r.get("error"),
            r.get("quality_pct"), r.get("last_run_at")]
           for r in (payload.get("by_layer") or [])])

    sheet("Runs",
          ["Test layer", "Started", "Status", "Passed", "Failed", "Checkpoints", "Pass %"],
          [[r.get("mapping_name"), r.get("started_at"), r.get("status"), r.get("pass_count"),
            r.get("fail_count"), r.get("total_checkpoints"), r.get("pass_pct")]
           for r in (payload.get("trend") or [])])

    sheet("Summary", ["Measure", "Value", "Detail"],
          [list(r) for r in headline_rows(payload)]
          + [["Scope", describe_scope(payload), ""]]
          + [["Note", n, ""] for n in caveats(payload)])

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def to_pbids(payload, _charts=None, api_base="http://127.0.0.1:5000"):
    """
    A Power BI Data Source file: opening it launches Power BI Desktop straight
    into Get Data against the live analytics endpoint, so the report refreshes
    rather than freezing at export time. Tiny JSON - it holds no data itself,
    only where to fetch it.
    """
    params = []
    if payload.get("scope_mapping_ids"):
        params.append("mapping_ids=" + ",".join(payload["scope_mapping_ids"]))
    params.append(f"basis={payload.get('basis', 'all')}")
    params.append(f"range={payload.get('range', 'all')}")
    url = f"{api_base}/api/s2d/analytics?" + "&".join(params)

    return json.dumps({
        "version": "0.1",
        "connections": [{
            "details": {"protocol": "http", "address": {"url": url}},
        }],
    }, indent=2).encode("utf-8")


BUILDERS = {
    "pdf": (to_pdf, "application/pdf", "pdf"),
    "docx": (to_docx,
             "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "docx"),
    "pptx": (to_pptx,
             "application/vnd.openxmlformats-officedocument.presentationml.presentation", "pptx"),
    "xlsx": (to_xlsx,
             "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "xlsx"),
    "pbids": (to_pbids, "application/json", "pbids"),
}
